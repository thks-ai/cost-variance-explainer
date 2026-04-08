from fastapi import FastAPI, UploadFile, File
from fastapi.responses import FileResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
import os
import sqlite3
import subprocess
from datetime import datetime
import shutil
import re

from pydantic import BaseModel
from sentence_transformers import SentenceTransformer
import numpy as np
import openpyxl
from docx import Document
import pdfplumber
from pptx import Presentation

# -----------------------------
# FastAPI アプリ
# -----------------------------
app = FastAPI()

app.mount("/static", StaticFiles(directory="static"), name="static")

# 既存ルーター
from app.api.rag import router as rag_router
from app.api.generate import router as generate_router
from app.api.upload import router as upload_router

# ★★★ 追加：analyze ルーター ★★★
from app.api.analyze import router as analyze_router

# 既存ルーター登録
app.include_router(rag_router)
app.include_router(generate_router)
app.include_router(upload_router)

# ★★★ 追加：/analyze を FastAPI に登録 ★★★
app.include_router(analyze_router)

# -----------------------------
# 画面
# -----------------------------
@app.get("/", response_class=HTMLResponse)
def index():
    return FileResponse("templates/index.html")

@app.get("/edit", response_class=HTMLResponse)
def edit_page():
    return FileResponse("templates/edit.html")

# ============================================================
# Excel 抽出（実務強化版）
# ============================================================
REASON_COL_KEYWORDS = [
    "理由", "原因", "要因", "コメント", "備考", "説明", "背景", "詳細"
]

def looks_like_reason(text: str) -> bool:
    if not text:
        return False
    if len(text) < 5:
        return False
    if re.fullmatch(r"[0-9\.\-]+", text):
        return False
    if len(re.sub(r"[0-9\.\-]", "", text)) == 0:
        return False
    return True

def extract_text_from_excel(filepath: str):
    texts = []
    try:
        wb = openpyxl.load_workbook(filepath, data_only=True)

        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if not cell:
                        continue
                    text = str(cell).strip()
                    if looks_like_reason(text):
                        texts.append(text)

        return list(set(texts))

    except Exception as e:
        return [f"Excel抽出エラー: {str(e)}"]

# ============================================================
# Word 抽出（箇条書き・表・テキストボックス対応）
# ============================================================
def extract_text_from_word(filepath: str):
    try:
        doc = Document(filepath)
        texts = []

        for p in doc.paragraphs:
            if p.text and p.text.strip():
                texts.append(p.text.strip())

            for run in p.runs:
                if run.text and run.text.strip():
                    texts.append(run.text.strip())

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        texts.append(cell_text)

        from docx.oxml import OxmlElement
        for el in doc.element.body.iter():
            if el.tag.endswith('}txbxContent'):
                for p in el.iter():
                    if p.tag.endswith('}p'):
                        text = ''.join([node.text for node in p.iter() if node.text])
                        if text and text.strip():
                            texts.append(text.strip())

        return list(set(texts))

    except Exception as e:
        return [f"Word抽出エラー: {str(e)}"]

# ============================================================
# PDF 抽出
# ============================================================
PDF_REASON_HINTS = [
    "理由", "原因", "要因", "ため", "影響", "結果", "改善", "問題", "課題"
]

def extract_text_from_pdf(filepath: str):
    texts = []
    try:
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                content = page.extract_text()
                if not content:
                    continue

                for line in content.split("\n"):
                    line = line.strip()
                    if not line:
                        continue

                    if re.match(r"^[・\-●■◆▶︎○◇\d\.\)]", line):
                        texts.append(line)
                        continue

                    if any(key in line for key in PDF_REASON_HINTS):
                        texts.append(line)
                        continue

                    if len(line) > 8 and re.search(r"[。.!?]$", line):
                        texts.append(line)

        return list(set(texts))

    except Exception as e:
        return [f"PDF抽出エラー: {str(e)}"]

# ============================================================
# PowerPoint 抽出
# ============================================================
def extract_text_from_ppt(filepath: str):
    try:
        prs = Presentation(filepath)
        texts = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    for line in shape.text.split("\n"):
                        line = line.strip()
                        if line:
                            texts.append(line)

                if hasattr(shape, "text_frame") and shape.text_frame:
                    for p in shape.text_frame.paragraphs:
                        if p.text and p.text.strip():
                            texts.append(p.text.strip())

        return list(set(texts))

    except Exception as e:
        return [f"PPT抽出エラー: {str(e)}"]

# ============================================================
# DB 登録
# ============================================================
clean_model = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")

def insert_into_db(text_list: list):
    cleaned = []
    for t in text_list:
        if not t:
            continue
        text = t.strip()
        if len(text) < 5:
            continue
        if text in cleaned:
            continue
        cleaned.append(text)

    if not cleaned:
        return

    embeddings = clean_model.encode(cleaned).astype("float32")

    final_texts = []
    for i, text in enumerate(cleaned):
        vec = embeddings[i]
        is_duplicate = False

        for existing in final_texts:
            existing_vec = embeddings[cleaned.index(existing)]
            sim = np.dot(vec, existing_vec) / (
                np.linalg.norm(vec) * np.linalg.norm(existing_vec)
            )
            if sim > 0.85:
                is_duplicate = True
                break

        if not is_duplicate:
            final_texts.append(text)

    conn = sqlite3.connect("data/reason_log.db")
    cur = conn.cursor()

    cur.execute("""
        CREATE TABLE IF NOT EXISTS reason_log (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            reason TEXT,
            tag TEXT DEFAULT 'その他'
        )
    """)

    for text in final_texts:
        cur.execute(
            "INSERT INTO reason_log (reason, tag) VALUES (?, ?)",
            (text, "その他")
        )

    conn.commit()
    conn.close()

# ============================================================
# ベクトル再構築
# ============================================================
def rebuild_vector_index():
    try:
        subprocess.run(
            ["python", "scripts/build_vector_index.py"],
            check=True
        )
        return "ベクトル再構築完了"
    except Exception as e:
        return f"ベクトル再構築エラー: {str(e)}"

# ============================================================
# 理由文一覧 API
# ============================================================
@app.get("/reasons")
def get_reasons():
    conn = sqlite3.connect("data/reason_log.db")
    cur = conn.cursor()

    cur.execute("SELECT id, reason, tag FROM reason_log ORDER BY id DESC")
    rows = cur.fetchall()

    conn.close()

    return [
        {"id": r[0], "reason": r[1], "tag": r[2] or "その他"}
        for r in rows
    ]

# ============================================================
# 理由文更新 API
# ============================================================
class ReasonUpdate(BaseModel):
    id: int
    reason: str
    tag: str

@app.post("/update_reason")
def update_reason(data: ReasonUpdate):
    conn = sqlite3.connect("data/reason_log.db")
    cur = conn.cursor()

    cur.execute(
        "UPDATE reason_log SET reason = ?, tag = ? WHERE id = ?",
        (data.reason, data.tag, data.id)
    )
    conn.commit()
    conn.close()

    rebuild_vector_index()

    return {"status": "ok", "message": "更新完了"}

# ============================================================
# ドラッグ＆ドロップ用：1ファイルずつ保存
# ============================================================
@app.post("/upload")
async def upload_single_file(file: UploadFile = File(...)):
    os.makedirs("uploaded", exist_ok=True)

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
    safe_filename = f"{timestamp}_{file.filename}"

    save_path = os.path.join("uploaded", safe_filename)

    with open(save_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    return {"status": "ok", "saved_as": safe_filename}

# ============================================================
# まとめて解析
# ============================================================
@app.post("/upload_files")
async def upload_files(files: list[UploadFile] = File(...)):
    base_dir = "data/uploads"
    os.makedirs(base_dir, exist_ok=True)

    results = []
    all_extracted_texts = []

    for file in files:
        ext = file.filename.lower().split(".")[-1]

        save_dir = f"{base_dir}/{ext}"
        os.makedirs(save_dir, exist_ok=True)

        filepath = f"{save_dir}/{file.filename}"
        contents = await file.read()

        with open(filepath, "wb") as f:
            f.write(contents)

        if ext in ["xlsx", "xls"]:
            extracted = extract_text_from_excel(filepath)
        elif ext == "docx":
            extracted = extract_text_from_word(filepath)
        elif ext == "pdf":
            extracted = extract_text_from_pdf(filepath)
        elif ext == "pptx":
            extracted = extract_text_from_ppt(filepath)
        else:
            extracted = [f"{file.filename} は未対応形式です"]

        all_extracted_texts.extend(extracted)

        results.append({
            "file": file.filename,
            "type": ext,
            "extracted_texts": extracted
        })

    insert_into_db(all_extracted_texts)

    vector_status = rebuild_vector_index()

    return {
        "status": "ok",
        "message": "ファイル抽出 → DB登録 → ベクトル再構築 完了",
        "vector_status": vector_status,
        "results": results
    }

# ============================================================
# 動作確認
# ============================================================
@app.get("/health")
def health():
    return {"status": "ok"}




















