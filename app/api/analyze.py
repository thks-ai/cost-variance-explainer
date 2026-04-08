from fastapi import APIRouter, UploadFile, File, Form
from typing import List
import pdfplumber
import docx
import openpyxl
from pptx import Presentation
from openai import OpenAI
import json
import io

from app.services.rag_search import rag_search

router = APIRouter()
client = OpenAI()

# ============================
# ファイル別テキスト抽出（BytesIO対応）
# ============================

def extract_text_from_pdf(file: UploadFile):
    contents = file.file.read()
    file.file.seek(0)
    text = ""
    with pdfplumber.open(io.BytesIO(contents)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text() or ""
            text += page_text + "\n"
    return text

def extract_text_from_docx(file: UploadFile):
    contents = file.file.read()
    file.file.seek(0)
    doc = docx.Document(io.BytesIO(contents))
    return "\n".join([p.text for p in doc.paragraphs if p.text])

def extract_text_from_excel(file: UploadFile):
    contents = file.file.read()
    file.file.seek(0)
    wb = openpyxl.load_workbook(io.BytesIO(contents), data_only=True)
    text = ""
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join([str(cell) for cell in row if cell is not None])
            if row_text.strip():
                text += row_text + "\n"
    return text

def extract_text_from_pptx(file: UploadFile):
    contents = file.file.read()
    file.file.seek(0)
    prs = Presentation(io.BytesIO(contents))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text += shape.text + "\n"
    return text

def extract_text_from_txt(file: UploadFile):
    return file.file.read().decode("utf-8")

def extract_text(file: UploadFile):
    name = file.filename.lower()
    if name.endswith(".pdf"):
        return extract_text_from_pdf(file)
    elif name.endswith(".docx"):
        return extract_text_from_docx(file)
    elif name.endswith(".xlsx"):
        return extract_text_from_excel(file)
    elif name.endswith(".pptx"):
        return extract_text_from_pptx(file)
    elif name.endswith(".txt"):
        return extract_text_from_txt(file)
    return ""

# ============================
# テキスト長制限
# ============================

def truncate(text: str, max_chars: int = 2000):
    return text if len(text) <= max_chars else text[:max_chars]

# ============================
# /analyze（ハイライト対応・高速版）
# ============================

@router.post("/analyze")
async def analyze_files(
    files: List[UploadFile] = File(...),
    question: str = Form(...),
    output_style: str = Form(...)
):
    merged_text = ""

    # ① 全ファイル抽出
    for file in files:
        extracted = extract_text(file)
        merged_text += f"\n--- {file.filename} ---\n{extracted}\n"

    # ② 長すぎる資料は圧縮
    merged_text = truncate(merged_text, 2000)

    # ③ RAG（高速化：top_k=3）
    rag_results = rag_search(question, top_k=3)
    rag_text = "\n".join([truncate(r["text"], 200) for r in rag_results])

    # ④ AI プロンプト（ハイライト対応）
    prompt = f"""
あなたは製造業の原価分析に詳しいアナリストです。

以下の資料内容と、過去の理由文データ（RAG検索結果）を読み込み、
ユーザーの質問に対して「最も妥当な原因」を生成してください。

【資料内容（要約用）】
{merged_text}

【RAG（過去データの候補）】
{rag_text}

【質問】
{question}

【出力形式】
{output_style}

【出力フォーマット（厳守）】
JSON形式で返すこと：

{{
  "answer": "回答文",
  "references": [
    {{
      "file": "ファイル名 or RAG",
      "text": "根拠となった文章（資料からそのまま抜粋）",
      "highlight": "回答に直接使った部分（text の中の exact match）"
    }}
  ]
}}

【要件】
- 資料に書かれている内容を最優先
- 資料に無い部分は RAG（過去データ）で補完
- highlight は text の中に必ず含まれる exact match とする
- JSON 以外のテキストは一切出力しない
"""

    res = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.2
    )

    raw = res.choices[0].message.content

    try:
        data = json.loads(raw)
    except:
        data = {"answer": raw, "references": []}

    return data






